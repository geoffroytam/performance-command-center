"""PowerPoint report generation -- Objectif Lune Command Center.

Generates a polished, template-matched PPTX with:
  - Dark section divider slides (Utopia Std Display / Georgia fallback)
  - White content slides with KPI cards, styled tables, and auto-insights
  - 16:9 widescreen (13.33 x 7.5 inches)
  - Objectif Lune brand palette

Every slide is built on a blank layout with manually positioned shapes
for full visual control.
"""

import pandas as pd
import numpy as np
from io import BytesIO
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

from utils.calculations import (
    calculate_baselines,
    aggregate_by_period,
    aggregate_for_date,
    format_currency,
    format_number,
    format_pct,
)
from utils.constants import COLORS, PLATFORM_COLORS, ROAS_TARGETS, load_settings


# =====================================================================
#  DESIGN TOKENS -- matched to user's actual PowerPoint template
# =====================================================================

# -- Fonts (primary + fallback) ----------------------------------------
FONT_TITLE = "Utopia Std Display"
FONT_TITLE_FALLBACK = "Georgia"
FONT_BODY = "Acumin Pro"
FONT_BODY_FALLBACK = "Calibri"

# -- Colours (RGBColor) ------------------------------------------------
DARK_BG = RGBColor(0x2D, 0x3E, 0x50)       # slide background for dividers / title
BODY_TEXT = RGBColor(0x40, 0x38, 0x33)       # #403833 primary body text
ACCENT_ORANGE = RGBColor(0xEB, 0x7C, 0x31)  # #EB7C31 accent borders / highlights
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xF7)

# -- Theme accent colours ----------------------------------------------
BLUE = RGBColor(0x4A, 0x6F, 0xA5)
GREEN = RGBColor(0x6B, 0x8F, 0x71)
RED = RGBColor(0xC4, 0x5C, 0x4A)
ORANGE = RGBColor(0xC7, 0x8B, 0x52)
GRAY = RGBColor(0x7A, 0x7A, 0x72)
LIGHT_GRAY = RGBColor(0xF0, 0xED, 0xE6)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)

# -- Table colours ------------------------------------------------------
TABLE_HEADER_BG = RGBColor(0x2D, 0x3E, 0x50)
TABLE_HEADER_FG = WHITE
TABLE_ROW_EVEN = RGBColor(0xF5, 0xF0, 0xE8)   # cream
TABLE_ROW_ODD = WHITE
TABLE_BORDER = RGBColor(0xD5, 0xD0, 0xC8)

# -- KPI card colours ---------------------------------------------------
KPI_CARD_BG = RGBColor(0xF5, 0xF0, 0xE8)
KPI_CARD_BORDER = RGBColor(0xEB, 0x7C, 0x31)

# -- Slide dimensions (16:9) -------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# -- Margins / gutters --------------------------------------------------
MARGIN_LEFT = Inches(0.8)
MARGIN_RIGHT = Inches(0.8)
CONTENT_W = SLIDE_W - MARGIN_LEFT - MARGIN_RIGHT  # ~11.73 in


# =====================================================================
#  XML namespace shorthand
# =====================================================================
_NSMAP = "http://schemas.openxmlformats.org/drawingml/2006/main"


# =====================================================================
#  Font helpers
# =====================================================================

def _apply_font(run, name, size, color, bold=False, italic=False):
    """Apply font properties to a single Run object."""
    run.font.name = name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def _set_paragraph_font(paragraph, name, size, color, bold=False,
                        alignment=PP_ALIGN.LEFT):
    """Set font on the first run of a paragraph (creates one if needed)."""
    paragraph.alignment = alignment
    if paragraph.runs:
        for run in paragraph.runs:
            _apply_font(run, name, size, color, bold)
    else:
        run = paragraph.add_run()
        run.text = paragraph.text if paragraph.text else ""
        _apply_font(run, name, size, color, bold)


# =====================================================================
#  Low-level drawing helpers
# =====================================================================

def _fill_slide_background(slide, color):
    """Set the entire slide background to a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_cell_border(cell, color=TABLE_BORDER, width=Pt(0.75)):
    """Apply a thin border to all four sides of a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    color_hex = f"{color.red:02X}{color.green:02X}{color.blue:02X}" if hasattr(color, "red") else "D5D0C8"
    for edge in ("lnL", "lnR", "lnT", "lnB"):
        ln = etree.SubElement(tcPr, f"{{{_NSMAP}}}{edge}")
        ln.set("w", str(int(width)))
        solidFill = etree.SubElement(ln, f"{{{_NSMAP}}}solidFill")
        srgbClr = etree.SubElement(solidFill, f"{{{_NSMAP}}}srgbClr")
        srgbClr.set("val", color_hex)


def _add_textbox(slide, left, top, width, height, text,
                 font_name=None, font_size=Pt(11), font_color=BODY_TEXT,
                 bold=False, alignment=PP_ALIGN.LEFT, word_wrap=True,
                 line_spacing=None):
    """Add a positioned text box and return its text frame.

    If font_name is None, uses FONT_BODY (Acumin Pro / Calibri fallback).
    """
    if font_name is None:
        font_name = FONT_BODY
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    if line_spacing is not None:
        p.line_spacing = line_spacing
    return tf


def _add_slide_title(slide, text, left=None, top=Inches(0.5),
                     font_size=Pt(28), font_color=BODY_TEXT):
    """Place a slide title in Utopia Std Display (Georgia fallback)."""
    if left is None:
        left = MARGIN_LEFT
    _add_textbox(
        slide, left, top, CONTENT_W, Inches(0.7),
        text, font_name=FONT_TITLE, font_size=font_size,
        font_color=font_color, bold=True, alignment=PP_ALIGN.LEFT,
    )


def _add_subtitle_line(slide, text, top=Inches(1.15)):
    """Place a subtitle line below the slide title (Acumin Pro)."""
    _add_textbox(
        slide, MARGIN_LEFT, top, CONTENT_W, Inches(0.4),
        text, font_name=FONT_BODY, font_size=Pt(13),
        font_color=GRAY, bold=False, alignment=PP_ALIGN.LEFT,
    )


def _add_accent_rule(slide, top, width=None, color=ACCENT_ORANGE):
    """Draw an orange accent line across the slide."""
    if width is None:
        width = CONTENT_W
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, top, width, Pt(3),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def _add_footer(slide, text="Objectif Lune -- Command Center",
                font_color=GRAY):
    """Add a subtle bottom-right footer."""
    _add_textbox(
        slide, SLIDE_W - Inches(4.5), SLIDE_H - Inches(0.55),
        Inches(4.0), Inches(0.35),
        text, font_name=FONT_BODY, font_size=Pt(8),
        font_color=font_color, bold=False, alignment=PP_ALIGN.RIGHT,
    )


# =====================================================================
#  KPI card drawing
# =====================================================================

def _draw_kpi_card(slide, left, top, width, height, label, value,
                   delta=None, value_color=BODY_TEXT):
    """Draw a single KPI card: cream rounded rectangle with orange left accent."""
    # Main card body
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    card.fill.solid()
    card.fill.fore_color.rgb = KPI_CARD_BG
    # Subtle border
    card.line.color.rgb = RGBColor(0xE0, 0xDC, 0xD5)
    card.line.width = Pt(0.75)
    # Reduce corner rounding
    card.adjustments[0] = 0.06

    # Orange accent strip on the left
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top + Inches(0.15), Pt(4), height - Inches(0.3),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_ORANGE
    accent.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.05)

    # Label
    p_label = tf.paragraphs[0]
    p_label.text = label.upper()
    p_label.font.name = FONT_BODY
    p_label.font.size = Pt(9)
    p_label.font.color.rgb = GRAY
    p_label.font.bold = True
    p_label.alignment = PP_ALIGN.CENTER
    p_label.space_after = Pt(2)

    # Value
    p_val = tf.add_paragraph()
    p_val.text = str(value)
    p_val.font.name = FONT_TITLE
    p_val.font.size = Pt(22)
    p_val.font.bold = True
    p_val.font.color.rgb = value_color
    p_val.alignment = PP_ALIGN.CENTER
    p_val.space_before = Pt(2)
    p_val.space_after = Pt(1)

    # Delta badge
    if delta is not None and delta != "":
        p_delta = tf.add_paragraph()
        delta_str = str(delta)
        is_positive = delta_str.startswith("+")
        is_negative = delta_str.startswith("-")
        if is_positive:
            arrow = "^ "
            d_color = GREEN
        elif is_negative:
            arrow = "v "
            d_color = RED
        else:
            arrow = ""
            d_color = GRAY
        p_delta.text = f"{arrow}{delta_str}"
        p_delta.font.name = FONT_BODY
        p_delta.font.size = Pt(9)
        p_delta.font.bold = True
        p_delta.font.color.rgb = d_color
        p_delta.alignment = PP_ALIGN.CENTER
        p_delta.space_before = Pt(0)


def _draw_kpi_row(slide, kpis, top=Inches(1.8), card_h=Inches(1.35)):
    """Draw a centred row of KPI cards with even spacing."""
    n = len(kpis)
    if n == 0:
        return
    card_w = Inches(2.6)
    gap = Inches(0.35)
    total_w = n * card_w + (n - 1) * gap
    start_left = int((SLIDE_W - total_w) / 2)

    for i, kpi in enumerate(kpis):
        left = start_left + int(i * (card_w + gap))
        _draw_kpi_card(
            slide, left, top, card_w, card_h,
            label=kpi["label"],
            value=kpi["value"],
            delta=kpi.get("delta"),
            value_color=kpi.get("color", BODY_TEXT),
        )


# =====================================================================
#  Table helper
# =====================================================================

def _add_styled_table(slide, headers, rows, top=Inches(3.3),
                      left=None, width=None, col_widths=None,
                      roas_col_idx=None, font_size=Pt(9)):
    """Add a table with dark header, alternating cream/white rows, and borders.

    If *roas_col_idx* is given, values in that column are colour-coded
    green/red vs the ROAS target for the row's campaign type.
    """
    if left is None:
        left = MARGIN_LEFT
    if width is None:
        width = CONTENT_W

    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    row_h = Inches(0.35)
    height = row_h * n_rows

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    else:
        # Equal distribution
        col_w = width // n_cols
        for i in range(n_cols):
            tbl.columns[i].width = col_w

    # ---- Header row ----
    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.name = FONT_BODY
        p.font.size = font_size
        p.font.bold = True
        p.font.color.rgb = TABLE_HEADER_FG
        p.alignment = PP_ALIGN.CENTER

        _set_cell_border(cell, color=RGBColor(0x1C, 0x2A, 0x3A), width=Pt(0.5))

    # ---- Data rows ----
    for r, row_data in enumerate(rows):
        bg = TABLE_ROW_EVEN if r % 2 == 0 else TABLE_ROW_ODD
        for c, val in enumerate(row_data):
            cell = tbl.cell(r + 1, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = cell.text_frame.paragraphs[0]
            p.text = str(val) if val is not None else "\u2014"
            p.font.name = FONT_BODY
            p.font.size = font_size
            p.font.color.rgb = BODY_TEXT
            # Left-align text columns (first col), centre numeric columns
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER

            _set_cell_border(cell, color=TABLE_BORDER, width=Pt(0.5))

            # Colour ROAS column
            if roas_col_idx is not None and c == roas_col_idx:
                _color_roas_cell(cell, val, row_data)

    return tbl_shape


def _color_roas_cell(cell, roas_val, row_data):
    """Colour a ROAS cell green/red based on the campaign-type target."""
    try:
        roas_num = float(str(roas_val).replace(",", ""))
    except (ValueError, TypeError):
        return

    # Determine target: check column 1 for campaign type, fallback to Prospecting target
    ctype = str(row_data[1]).strip() if len(row_data) > 1 else ""
    target = ROAS_TARGETS.get(ctype, ROAS_TARGETS.get("Prospecting", 8))
    color = GREEN if roas_num >= target else RED
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = color
        p.font.bold = True


# =====================================================================
#  Takeaway generator
# =====================================================================

def _generate_takeaways(filtered: pd.DataFrame) -> list[str]:
    """Return 4-6 auto-generated bullet points from the data."""
    bullets = []

    total_spend = filtered["spend"].sum()
    total_revenue = filtered["revenue"].sum()
    total_orders = filtered["conversions"].sum()
    total_impressions = filtered["impressions"].sum()
    total_clicks = filtered["clicks"].sum()
    overall_roas = total_revenue / total_spend if total_spend > 0 else 0
    overall_aov = total_revenue / total_orders if total_orders > 0 else 0
    overall_cpm = total_spend / total_impressions * 1000 if total_impressions > 0 else 0
    overall_ctr = total_clicks / total_impressions * 100 if total_impressions > 0 else 0
    overall_cvr = total_orders / total_clicks * 100 if total_clicks > 0 else 0

    plat_agg = (
        filtered.groupby("platform")
        .agg(
            spend=("spend", "sum"),
            revenue=("revenue", "sum"),
            conversions=("conversions", "sum"),
            impressions=("impressions", "sum"),
            clicks=("clicks", "sum"),
        )
        .reset_index()
    )
    plat_agg["roas"] = np.where(
        plat_agg["spend"] > 0, plat_agg["revenue"] / plat_agg["spend"], 0
    )
    plat_agg["aov"] = np.where(
        plat_agg["conversions"] > 0, plat_agg["revenue"] / plat_agg["conversions"], 0
    )

    # 1. Best and worst platform by ROAS
    if not plat_agg.empty:
        best = plat_agg.loc[plat_agg["roas"].idxmax()]
        worst = plat_agg.loc[plat_agg["roas"].idxmin()]
        bullets.append(
            f"Best performing platform: {best['platform']} with a ROAS of "
            f"{best['roas']:.1f}x on {format_currency(best['spend'])} spend, "
            f"generating {format_currency(best['revenue'])} in revenue."
        )
        if best["platform"] != worst["platform"]:
            bullets.append(
                f"Lowest performing platform: {worst['platform']} with a ROAS of "
                f"{worst['roas']:.1f}x ({format_currency(worst['spend'])} spend)."
            )

    # 2. Overall performance
    bullets.append(
        f"Overall blended ROAS: {overall_roas:.1f}x across {format_currency(total_spend)} "
        f"total spend with {format_number(total_orders, 0)} orders "
        f"(AOV: {format_currency(overall_aov)})."
    )

    # 3. Efficiency metrics
    bullets.append(
        f"Funnel efficiency: CPM {format_currency(overall_cpm)}, "
        f"CTR {overall_ctr:.2f}%, CVR {overall_cvr:.2f}%."
    )

    # 4. Check platforms vs ROAS targets
    prosp_target = ROAS_TARGETS.get("Prospecting", 8)
    retarg_target = ROAS_TARGETS.get("Retargeting", 14)

    exceeding = []
    below = []
    for _, row in plat_agg.iterrows():
        plat = row["platform"]
        roas = row["roas"]
        if roas >= prosp_target:
            exceeding.append(f"{plat} ({roas:.1f}x)")
        elif roas < prosp_target * 0.7:
            below.append(f"{plat} ({roas:.1f}x)")

    if exceeding:
        bullets.append(
            f"Meeting/exceeding Prospecting target ({prosp_target}x): "
            + ", ".join(exceeding) + "."
        )
    if below:
        bullets.append(
            f"Below target -- review spend allocation: "
            + ", ".join(below)
            + f" (target: {prosp_target}x)."
        )

    # 5. Recommendation
    if overall_roas >= prosp_target:
        bullets.append(
            "Recommendation: Performance is strong. Consider scaling spend on top "
            "performers while monitoring CPM inflation."
        )
    else:
        bullets.append(
            "Recommendation: Blended ROAS is below target. Evaluate creative fatigue, "
            "audience saturation, and consider reallocating budget from underperformers."
        )

    return bullets[:6]


# =====================================================================
#  Prior-period comparison helper
# =====================================================================

def _compute_prior_period_kpis(df, start_date, end_date, platforms=None):
    """Return KPIs dict for the prior period of equal length, or None."""
    start = pd.Timestamp(start_date)
    end = pd.Timestamp(end_date)
    duration = (end - start).days
    prior_end = start - timedelta(days=1)
    prior_start = prior_end - timedelta(days=duration)

    mask = df["date"].between(prior_start, prior_end)
    if platforms:
        mask = mask & df["platform"].isin(platforms)
    prior = df[mask]

    if prior.empty:
        return None

    spend = prior["spend"].sum()
    revenue = prior["revenue"].sum()
    roas = revenue / spend if spend > 0 else 0
    orders = prior["conversions"].sum()
    impressions = prior["impressions"].sum()
    clicks = prior["clicks"].sum()
    cpm = spend / impressions * 1000 if impressions > 0 else 0
    ctr = clicks / impressions * 100 if impressions > 0 else 0
    cvr = orders / clicks * 100 if clicks > 0 else 0
    aov = revenue / orders if orders > 0 else 0

    return {
        "spend": spend,
        "revenue": revenue,
        "roas": roas,
        "orders": orders,
        "impressions": impressions,
        "clicks": clicks,
        "cpm": cpm,
        "ctr": ctr,
        "cvr": cvr,
        "aov": aov,
    }


def _delta_str(current, prior):
    """Return a formatted +/-% string, or empty string if unavailable."""
    if prior is None or prior == 0:
        return ""
    pct = (current - prior) / abs(prior) * 100
    sign = "+" if pct >= 0 else ""
    return f"{sign}{pct:.1f}%"


# =====================================================================
#  SLIDE BUILDERS
# =====================================================================

# -- 1. Title Slide ----------------------------------------------------

def _build_title_slide(prs, date_range_str):
    """Slide 1 -- Dark background title page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _fill_slide_background(slide, DARK_BG)

    # Subtle top accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_W, Pt(5),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_ORANGE
    accent.line.fill.background()

    # Main title
    _add_textbox(
        slide, MARGIN_LEFT, Inches(2.2), CONTENT_W, Inches(1.2),
        "Performance Report",
        font_name=FONT_TITLE, font_size=Pt(48), font_color=WHITE,
        bold=True, alignment=PP_ALIGN.CENTER,
    )

    # Date range subtitle
    _add_textbox(
        slide, MARGIN_LEFT, Inches(3.5), CONTENT_W, Inches(0.6),
        date_range_str,
        font_name=FONT_BODY, font_size=Pt(18), font_color=RGBColor(0xBB, 0xBB, 0xBB),
        bold=False, alignment=PP_ALIGN.CENTER,
    )

    # Orange separator
    sep_w = Inches(2.0)
    sep_left = int((SLIDE_W - sep_w) / 2)
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        sep_left, Inches(4.3), sep_w, Pt(3),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = ACCENT_ORANGE
    sep.line.fill.background()

    # Branding
    _add_textbox(
        slide, MARGIN_LEFT, Inches(5.8), CONTENT_W, Inches(0.5),
        "Objectif Lune",
        font_name=FONT_TITLE, font_size=Pt(20), font_color=ACCENT_ORANGE,
        bold=True, alignment=PP_ALIGN.CENTER,
    )

    _add_textbox(
        slide, MARGIN_LEFT, Inches(6.3), CONTENT_W, Inches(0.4),
        "Command Center",
        font_name=FONT_BODY, font_size=Pt(12),
        font_color=RGBColor(0x99, 0x99, 0x99),
        bold=False, alignment=PP_ALIGN.CENTER,
    )

    return slide


# -- 2. Executive Summary -----------------------------------------------

def _build_executive_summary(prs, filtered, prior_kpis):
    """Slide 2 -- Executive Summary with 4 KPI cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, WHITE)

    _add_slide_title(slide, "Executive Summary")
    _add_accent_rule(slide, Inches(1.15))

    total_spend = filtered["spend"].sum()
    total_revenue = filtered["revenue"].sum()
    total_orders = filtered["conversions"].sum()
    blended_roas = total_revenue / total_spend if total_spend > 0 else 0

    # Delta strings
    d_roas = _delta_str(blended_roas, prior_kpis["roas"]) if prior_kpis else ""
    d_spend = _delta_str(total_spend, prior_kpis["spend"]) if prior_kpis else ""
    d_rev = _delta_str(total_revenue, prior_kpis["revenue"]) if prior_kpis else ""
    d_orders = _delta_str(total_orders, prior_kpis["orders"]) if prior_kpis else ""

    # ROAS colour coding
    prosp_target = ROAS_TARGETS.get("Prospecting", 8)
    roas_color = GREEN if blended_roas >= prosp_target else RED

    kpis = [
        {"label": "Blended ROAS", "value": f"{blended_roas:.1f}x", "delta": d_roas, "color": roas_color},
        {"label": "Total Spend", "value": format_currency(total_spend), "delta": d_spend},
        {"label": "Total Revenue", "value": format_currency(total_revenue), "delta": d_rev},
        {"label": "Total Orders", "value": format_number(total_orders, 0), "delta": d_orders},
    ]
    _draw_kpi_row(slide, kpis, top=Inches(1.7), card_h=Inches(1.45))

    # Summary sentence below KPIs
    total_impressions = filtered["impressions"].sum()
    total_clicks = filtered["clicks"].sum()
    cpm = total_spend / total_impressions * 1000 if total_impressions > 0 else 0
    ctr = total_clicks / total_impressions * 100 if total_impressions > 0 else 0
    cvr = total_orders / total_clicks * 100 if total_clicks > 0 else 0
    aov = total_revenue / total_orders if total_orders > 0 else 0

    summary_text = (
        f"CPM: {format_currency(cpm)}  |  CTR: {ctr:.2f}%  |  "
        f"CVR: {cvr:.2f}%  |  AOV: {format_currency(aov)}"
    )
    _add_textbox(
        slide, MARGIN_LEFT, Inches(3.5), CONTENT_W, Inches(0.4),
        summary_text,
        font_name=FONT_BODY, font_size=Pt(12), font_color=GRAY,
        bold=False, alignment=PP_ALIGN.CENTER,
    )

    # Period comparison note
    if prior_kpis:
        _add_textbox(
            slide, MARGIN_LEFT, Inches(3.95), CONTENT_W, Inches(0.35),
            "Deltas shown vs. prior period of equal length",
            font_name=FONT_BODY, font_size=Pt(9), font_color=GRAY,
            bold=False, alignment=PP_ALIGN.CENTER,
        )

    # Additional detail: mini per-platform summary at the bottom
    plat_summary = (
        filtered.groupby("platform")
        .agg(spend=("spend", "sum"), revenue=("revenue", "sum"),
             conversions=("conversions", "sum"))
        .reset_index()
    )
    plat_summary["roas"] = np.where(
        plat_summary["spend"] > 0, plat_summary["revenue"] / plat_summary["spend"], 0
    )
    plat_summary = plat_summary.sort_values("revenue", ascending=False)

    if not plat_summary.empty:
        headers = ["Platform", "Spend", "Revenue", "ROAS", "Orders"]
        rows = []
        for _, r in plat_summary.iterrows():
            rows.append([
                r["platform"],
                format_currency(r["spend"]),
                format_currency(r["revenue"]),
                f"{r['roas']:.1f}",
                format_number(r["conversions"], 0),
            ])
        _add_styled_table(
            slide, headers, rows,
            top=Inches(4.5),
            col_widths=[2.5, 2.5, 2.5, 2.0, 2.23],
            roas_col_idx=3,
            font_size=Pt(10),
        )

    _add_footer(slide)
    return slide


# -- 3. Section Divider ------------------------------------------------

def _build_section_divider(prs, title, subtitle=""):
    """Dark background section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, DARK_BG)

    # Top accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_W, Pt(5),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_ORANGE
    accent.line.fill.background()

    # Centred title
    _add_textbox(
        slide, MARGIN_LEFT, Inches(2.6), CONTENT_W, Inches(1.2),
        title,
        font_name=FONT_TITLE, font_size=Pt(48), font_color=WHITE,
        bold=True, alignment=PP_ALIGN.CENTER,
    )

    # Orange underline
    line_w = Inches(3.0)
    line_left = int((SLIDE_W - line_w) / 2)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        line_left, Inches(3.9), line_w, Pt(3),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()

    # Subtitle
    if subtitle:
        _add_textbox(
            slide, MARGIN_LEFT, Inches(4.2), CONTENT_W, Inches(0.6),
            subtitle,
            font_name=FONT_BODY, font_size=Pt(16),
            font_color=RGBColor(0xBB, 0xBB, 0xBB),
            bold=False, alignment=PP_ALIGN.CENTER,
        )

    return slide


# -- 4. Platform Breakdown Table ----------------------------------------

def _build_platform_table_slide(prs, filtered):
    """Platform Breakdown -- styled table with all key metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, WHITE)

    _add_slide_title(slide, "Platform Breakdown")
    _add_subtitle_line(slide, "Performance by platform and campaign type")
    _add_accent_rule(slide, Inches(1.55))

    plat_summary = (
        filtered.groupby(["platform", "campaign_type"])
        .agg(
            spend=("spend", "sum"),
            revenue=("revenue", "sum"),
            impressions=("impressions", "sum"),
            clicks=("clicks", "sum"),
            conversions=("conversions", "sum"),
        )
        .reset_index()
    )
    plat_summary["roas"] = np.where(
        plat_summary["spend"] > 0,
        plat_summary["revenue"] / plat_summary["spend"], 0,
    )
    plat_summary["cpm"] = np.where(
        plat_summary["impressions"] > 0,
        plat_summary["spend"] / plat_summary["impressions"] * 1000, 0,
    )
    plat_summary["ctr"] = np.where(
        plat_summary["impressions"] > 0,
        plat_summary["clicks"] / plat_summary["impressions"] * 100, 0,
    )
    plat_summary["cvr"] = np.where(
        plat_summary["clicks"] > 0,
        plat_summary["conversions"] / plat_summary["clicks"] * 100, 0,
    )
    plat_summary["aov"] = np.where(
        plat_summary["conversions"] > 0,
        plat_summary["revenue"] / plat_summary["conversions"], 0,
    )
    plat_summary = plat_summary.sort_values(["platform", "campaign_type"])

    headers = ["Platform", "Type", "Spend", "Revenue", "ROAS",
               "Orders", "AOV", "CPM", "CTR", "CVR"]
    rows = []
    for _, r in plat_summary.iterrows():
        rows.append([
            r["platform"],
            r["campaign_type"],
            format_currency(r["spend"]),
            format_currency(r["revenue"]),
            f"{r['roas']:.1f}",
            format_number(r["conversions"], 0),
            format_currency(r["aov"]),
            format_currency(r["cpm"]),
            f"{r['ctr']:.2f}%",
            f"{r['cvr']:.2f}%",
        ])

    _add_styled_table(
        slide, headers, rows,
        top=Inches(1.9),
        col_widths=[1.3, 1.1, 1.2, 1.2, 0.85, 0.85, 1.1, 1.1, 0.95, 0.95],
        roas_col_idx=4,
        font_size=Pt(9),
    )

    _add_footer(slide)
    return slide


# -- 5. Per-Platform Detail Slide ----------------------------------------

def _build_platform_detail_slide(prs, filtered, df_full, platform,
                                 start_date, end_date, platforms_list):
    """One slide per platform with KPI cards + P/R table + funnel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, WHITE)

    _add_slide_title(slide, f"{platform}")
    _add_subtitle_line(slide, "Platform Deep Dive")
    _add_accent_rule(slide, Inches(1.55))

    plat_data = filtered[filtered["platform"] == platform]
    plat_spend = plat_data["spend"].sum()
    plat_revenue = plat_data["revenue"].sum()
    plat_roas = plat_revenue / plat_spend if plat_spend > 0 else 0
    plat_orders = plat_data["conversions"].sum()
    plat_impressions = plat_data["impressions"].sum()
    plat_clicks = plat_data["clicks"].sum()
    plat_cpm = plat_spend / plat_impressions * 1000 if plat_impressions > 0 else 0
    plat_ctr = plat_clicks / plat_impressions * 100 if plat_impressions > 0 else 0
    plat_cvr = plat_orders / plat_clicks * 100 if plat_clicks > 0 else 0
    plat_aov = plat_revenue / plat_orders if plat_orders > 0 else 0

    # Prior period for this platform
    prior = _compute_prior_period_kpis(df_full, start_date, end_date, [platform])
    d_roas = _delta_str(plat_roas, prior["roas"]) if prior else ""
    d_spend = _delta_str(plat_spend, prior["spend"]) if prior else ""
    d_rev = _delta_str(plat_revenue, prior["revenue"]) if prior else ""
    d_orders = _delta_str(plat_orders, prior["orders"]) if prior else ""

    # ROAS colour
    prosp_target = ROAS_TARGETS.get("Prospecting", 8)
    roas_color = GREEN if plat_roas >= prosp_target else RED

    kpis = [
        {"label": "ROAS", "value": f"{plat_roas:.1f}x", "delta": d_roas, "color": roas_color},
        {"label": "Spend", "value": format_currency(plat_spend), "delta": d_spend},
        {"label": "Revenue", "value": format_currency(plat_revenue), "delta": d_rev},
        {"label": "Orders", "value": format_number(plat_orders, 0), "delta": d_orders},
    ]
    _draw_kpi_row(slide, kpis, top=Inches(1.85), card_h=Inches(1.35))

    # -- Prospecting / Retargeting Split Table --
    type_agg = (
        plat_data.groupby("campaign_type")
        .agg(
            spend=("spend", "sum"),
            revenue=("revenue", "sum"),
            impressions=("impressions", "sum"),
            clicks=("clicks", "sum"),
            conversions=("conversions", "sum"),
        )
        .reset_index()
    )
    type_agg["roas"] = np.where(
        type_agg["spend"] > 0, type_agg["revenue"] / type_agg["spend"], 0,
    )
    type_agg["cpm"] = np.where(
        type_agg["impressions"] > 0,
        type_agg["spend"] / type_agg["impressions"] * 1000, 0,
    )
    type_agg["ctr"] = np.where(
        type_agg["impressions"] > 0,
        type_agg["clicks"] / type_agg["impressions"] * 100, 0,
    )
    type_agg["cvr"] = np.where(
        type_agg["clicks"] > 0,
        type_agg["conversions"] / type_agg["clicks"] * 100, 0,
    )
    type_agg["aov"] = np.where(
        type_agg["conversions"] > 0,
        type_agg["revenue"] / type_agg["conversions"], 0,
    )

    # Left side: P/R table
    pr_headers = ["Type", "Spend", "Revenue", "ROAS", "Orders", "CPM"]
    pr_rows = []
    for _, r in type_agg.iterrows():
        pr_rows.append([
            r["campaign_type"],
            format_currency(r["spend"]),
            format_currency(r["revenue"]),
            f"{r['roas']:.1f}",
            format_number(r["conversions"], 0),
            format_currency(r["cpm"]),
        ])

    table_left = MARGIN_LEFT
    table_w = Inches(7.0)
    _add_textbox(
        slide, table_left, Inches(3.5), table_w, Inches(0.35),
        "Prospecting / Retargeting Split",
        font_name=FONT_BODY, font_size=Pt(11), font_color=BODY_TEXT,
        bold=True, alignment=PP_ALIGN.LEFT,
    )

    tbl_shape = _add_styled_table(
        slide, pr_headers, pr_rows,
        top=Inches(3.9),
        left=table_left,
        width=table_w,
        col_widths=[1.3, 1.2, 1.2, 0.8, 0.8, 1.0],
        roas_col_idx=3,
        font_size=Pt(9),
    )

    # Right side: Funnel mini-cards
    funnel_left = Inches(8.3)
    funnel_w = Inches(4.2)

    _add_textbox(
        slide, funnel_left, Inches(3.5), funnel_w, Inches(0.35),
        "Funnel Metrics",
        font_name=FONT_BODY, font_size=Pt(11), font_color=BODY_TEXT,
        bold=True, alignment=PP_ALIGN.LEFT,
    )

    funnel_metrics = [
        ("CPM", format_currency(plat_cpm)),
        ("CTR", f"{plat_ctr:.2f}%"),
        ("CVR", f"{plat_cvr:.2f}%"),
        ("AOV", format_currency(plat_aov)),
    ]

    funnel_card_h = Inches(0.65)
    funnel_card_w = Inches(1.9)
    funnel_gap = Inches(0.15)

    for i, (label, value) in enumerate(funnel_metrics):
        row_idx = i // 2
        col_idx = i % 2
        f_left = funnel_left + col_idx * (funnel_card_w + funnel_gap)
        f_top = Inches(3.95) + row_idx * (funnel_card_h + funnel_gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            f_left, f_top, funnel_card_w, funnel_card_h,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = RGBColor(0xE0, 0xDC, 0xD5)
        card.line.width = Pt(0.5)
        card.adjustments[0] = 0.08

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_top = Inches(0.05)

        p_lbl = tf.paragraphs[0]
        p_lbl.text = label
        p_lbl.font.name = FONT_BODY
        p_lbl.font.size = Pt(8)
        p_lbl.font.color.rgb = GRAY
        p_lbl.font.bold = True
        p_lbl.alignment = PP_ALIGN.CENTER

        p_val = tf.add_paragraph()
        p_val.text = value
        p_val.font.name = FONT_TITLE
        p_val.font.size = Pt(16)
        p_val.font.bold = True
        p_val.font.color.rgb = BODY_TEXT
        p_val.alignment = PP_ALIGN.CENTER
        p_val.space_before = Pt(1)

    _add_footer(slide)
    return slide


# -- 6. Key Takeaways Slide ---------------------------------------------

def _build_takeaways_slide(prs, filtered):
    """Key Takeaways -- auto-generated bullet insights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, WHITE)

    _add_slide_title(slide, "Key Takeaways")
    _add_subtitle_line(slide, "Auto-generated insights from your data")
    _add_accent_rule(slide, Inches(1.55))

    bullets = _generate_takeaways(filtered)

    top = Inches(2.0)
    for i, bullet in enumerate(bullets):
        y = top + Inches(i * 0.78)

        # Orange bullet marker
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            MARGIN_LEFT + Inches(0.15), y + Inches(0.08),
            Inches(0.12), Inches(0.12),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = ACCENT_ORANGE
        marker.line.fill.background()

        # Bullet text
        _add_textbox(
            slide, MARGIN_LEFT + Inches(0.5), y, CONTENT_W - Inches(0.8), Inches(0.7),
            bullet,
            font_name=FONT_BODY, font_size=Pt(13), font_color=BODY_TEXT,
            bold=False, alignment=PP_ALIGN.LEFT,
        )

    _add_footer(slide)
    return slide


# -- 7. Closing Slide ---------------------------------------------------

def _build_closing_slide(prs):
    """Dark background closing/branding slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, DARK_BG)

    # Top accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_W, Pt(5),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_ORANGE
    accent.line.fill.background()

    # Brand name
    _add_textbox(
        slide, MARGIN_LEFT, Inches(2.5), CONTENT_W, Inches(1.0),
        "Objectif Lune",
        font_name=FONT_TITLE, font_size=Pt(44), font_color=WHITE,
        bold=True, alignment=PP_ALIGN.CENTER,
    )

    # Orange separator
    sep_w = Inches(2.5)
    sep_left = int((SLIDE_W - sep_w) / 2)
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        sep_left, Inches(3.6), sep_w, Pt(3),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = ACCENT_ORANGE
    sep.line.fill.background()

    # Subtitle
    _add_textbox(
        slide, MARGIN_LEFT, Inches(3.9), CONTENT_W, Inches(0.6),
        "Command Center",
        font_name=FONT_BODY, font_size=Pt(18),
        font_color=RGBColor(0xBB, 0xBB, 0xBB),
        bold=False, alignment=PP_ALIGN.CENTER,
    )

    # Tagline
    _add_textbox(
        slide, MARGIN_LEFT, Inches(5.0), CONTENT_W, Inches(0.5),
        "Data-driven performance, beautifully presented.",
        font_name=FONT_BODY, font_size=Pt(12),
        font_color=RGBColor(0x88, 0x88, 0x88),
        bold=False, alignment=PP_ALIGN.CENTER,
    )

    return slide


# =====================================================================
#  PUBLIC API
# =====================================================================

def generate_pptx_report(
    df: pd.DataFrame,
    start_date,
    end_date,
    platforms: list = None,
) -> BytesIO:
    """Generate a branded PowerPoint report and return a BytesIO buffer.

    Parameters
    ----------
    df : pd.DataFrame
        Full dataset with columns: date, platform, campaign_type, spend,
        revenue, impressions, clicks, conversions, etc.
    start_date, end_date : date-like
        Reporting window boundaries.
    platforms : list[str] | None
        Optional filter to limit which platforms are included.

    Returns
    -------
    BytesIO
        In-memory PPTX file ready for Streamlit download.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    settings = load_settings()

    # ── Filter data ──────────────────────────────────────────────
    mask = df["date"].between(pd.Timestamp(start_date), pd.Timestamp(end_date))
    if platforms:
        mask = mask & df["platform"].isin(platforms)
    filtered = df[mask].copy()

    # Date range string
    if hasattr(start_date, "strftime"):
        date_range_str = (
            f"{start_date.strftime('%d/%m/%Y')}  --  {end_date.strftime('%d/%m/%Y')}"
        )
    else:
        date_range_str = f"{start_date}  --  {end_date}"

    # Handle empty data
    if filtered.empty:
        _build_title_slide(prs, date_range_str)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _fill_slide_background(slide, WHITE)
        _add_slide_title(slide, "No Data Available")
        _add_textbox(
            slide, MARGIN_LEFT, Inches(2.5), CONTENT_W, Inches(1.0),
            "No data found for the selected period and platforms.",
            font_name=FONT_BODY, font_size=Pt(16), font_color=GRAY,
            alignment=PP_ALIGN.CENTER,
        )
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    # Prior period for comparison
    prior_kpis = _compute_prior_period_kpis(df, start_date, end_date, platforms)

    # ── Build slides ─────────────────────────────────────────────

    # 1. Title Page
    _build_title_slide(prs, date_range_str)

    # 2. Executive Summary
    _build_executive_summary(prs, filtered, prior_kpis)

    # 3. Section Divider: Platform Performance
    _build_section_divider(prs, "Platform Performance",
                           "Breakdown by platform and campaign type")

    # 4. Platform Breakdown Table
    _build_platform_table_slide(prs, filtered)

    # 5. Per-Platform Detail Slides
    for platform in sorted(filtered["platform"].unique()):
        _build_platform_detail_slide(
            prs, filtered, df, platform,
            start_date, end_date, platforms,
        )

    # 6. Section Divider: Key Takeaways
    _build_section_divider(prs, "Key Takeaways",
                           "Auto-generated insights and recommendations")

    # 7. Key Takeaways
    _build_takeaways_slide(prs, filtered)

    # 8. Closing Slide
    _build_closing_slide(prs)

    # ── Save ─────────────────────────────────────────────────────
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
